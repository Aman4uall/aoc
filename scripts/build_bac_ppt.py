from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/ayzaman/.gemini/antigravity/scratch/AoC")
ASSET_DIR = ROOT / "outputs/benzalkonium-chloride-benchmark-live-v26/presentation_assets"
OUT_DIR = ROOT / "outputs/bac-ppt"
OUT_PATH = OUT_DIR / "BAC_Plant_Design_Feasibility_v1.pptx"


COLORS = {
    "navy": RGBColor(28, 48, 84),
    "teal": RGBColor(11, 107, 87),
    "sand": RGBColor(245, 239, 232),
    "charcoal": RGBColor(39, 43, 52),
    "slate": RGBColor(95, 106, 126),
    "line": RGBColor(203, 210, 221),
    "accent": RGBColor(217, 163, 58),
    "danger": RGBColor(177, 63, 63),
    "warn": RGBColor(200, 122, 26),
    "bg2": RGBColor(251, 251, 252),
    "white": RGBColor(255, 255, 255),
}


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS["navy"]
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLORS["slate"]
    return box


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(12.4), Inches(6.9), Inches(0.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = str(idx)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["slate"]


def add_bullets(slide, bullets, x, y, w, h, font_size=18, color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS["charcoal"]
        p.space_after = Pt(5)
    return box


def add_table(slide, rows, cols, data, x, y, w, h, col_widths=None, header_fill=None, font_size=12):
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, width in enumerate(col_widths):
            tbl.columns[i].width = width
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["charcoal"]
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLORS["navy"]
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_fill or RGBColor(232, 237, 245)
            else:
                cell.fill.fore_color.rgb = COLORS["white"]
    return tbl


def add_note_box(slide, text, x, y, w, h, fill=None, font_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or RGBColor(245, 247, 250)
    shape.line.color.rgb = COLORS["line"]
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLORS["charcoal"]
    return shape


def add_image_contain(slide, image_path, x, y, w, h):
    pic = slide.shapes.add_picture(str(image_path), x, y, width=w)
    if pic.height > h:
        ratio = h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = x + int((w - pic.width) / 2)
    pic.top = y + int((h - pic.height) / 2)
    return pic


def add_metric_cards(slide, items, x, y, card_w, card_h, gap=0.18):
    for i, (label, value) in enumerate(items):
        left = x + Inches(i) * 0  # placeholder, handled below
        left = x + int(i * (card_w + Inches(gap)))
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["bg2"]
        shape.line.color.rgb = COLORS["line"]
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLORS["slate"]
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLORS["navy"]


def add_zone_diagram(slide):
    labels = [
        ("Tank Farm / Receipt", COLORS["accent"]),
        ("Reaction Zone", COLORS["danger"]),
        ("Separation & Finishing", COLORS["teal"]),
        ("Waste / Recovery", RGBColor(120, 120, 120)),
    ]
    x = Inches(0.8)
    y = Inches(1.7)
    w = Inches(2.65)
    h = Inches(1.3)
    gap = Inches(0.3)
    for i, (label, color) in enumerate(labels):
        left = x + i * (w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            min(color[0] + 35, 255), min(color[1] + 35, 255), min(color[2] + 35, 255)
        )
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS["charcoal"]
        if i < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left + w + Inches(0.05), y + Inches(0.45), Inches(0.18), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["slate"]
            arrow.line.color.rgb = COLORS["slate"]
    add_note_box(
        slide,
        "Dispatch edge and emergency access are kept away from the core reaction block, while utility and maintenance access are preserved across the main process train.",
        Inches(0.9),
        Inches(4.1),
        Inches(11.0),
        Inches(0.95),
        fill=RGBColor(244, 247, 250),
        font_size=13,
    )


def add_shew_diagram(slide):
    titles = [
        ("Air & Vapor Control", ["Closed handling", "Vent-management logic", "Storage vapor control"], COLORS["navy"]),
        ("Wastewater Handling", ["Segregated collection", "Neutralization / equalization", "Dahej CETP route"], COLORS["teal"]),
        ("Hazardous Waste", ["Off-spec liquid waste", "Contained solids storage", "Authorized disposal"], COLORS["warn"]),
    ]
    x = Inches(0.75)
    y = Inches(1.6)
    w = Inches(3.8)
    h = Inches(2.0)
    gap = Inches(0.3)
    for i, (title, bullets, color) in enumerate(titles):
        left = x + i * (w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["bg2"]
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color
        for bullet in bullets:
            pb = tf.add_paragraph()
            pb.text = bullet
            pb.font.size = Pt(13)
            pb.font.color.rgb = COLORS["charcoal"]
            pb.level = 0
    add_note_box(
        slide,
        "Bunded / diked storage, no direct discharge, and hazardous-waste segregation are retained as core SHE anchors.",
        Inches(0.9),
        Inches(4.15),
        Inches(11.0),
        Inches(0.75),
        fill=RGBColor(250, 244, 232),
        font_size=13,
    )


def add_utility_diagram(slide):
    process_blocks = ["Storage / Feed Prep", "Reactor", "Purification", "Storage / Waste"]
    utilities = [
        ("Steam", COLORS["accent"]),
        ("Cooling Water", RGBColor(90, 145, 210)),
        ("Power", COLORS["navy"]),
        ("Nitrogen", RGBColor(115, 115, 115)),
    ]
    x = Inches(1.1)
    y = Inches(2.0)
    w = Inches(2.2)
    h = Inches(0.9)
    gap = Inches(0.35)
    for i, label in enumerate(process_blocks):
        left = x + i * (w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(247, 249, 252)
        shape.line.color.rgb = COLORS["line"]
        p = shape.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLORS["charcoal"]
    rack = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(1.15), Inches(10.8), Inches(0.3))
    rack.fill.solid()
    rack.fill.fore_color.rgb = RGBColor(235, 239, 245)
    rack.line.color.rgb = COLORS["line"]
    for i, (u, color) in enumerate(utilities):
        tx = slide.shapes.add_textbox(Inches(1.0 + 2.6 * i), Inches(0.72), Inches(2.3), Inches(0.25))
        p = tx.text_frame.paragraphs[0]
        p.text = u
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
    for i in range(4):
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(1.75 + 2.55 * i), Inches(1.45), Inches(0.18), Inches(0.35))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["slate"]
        line.line.color.rgb = COLORS["slate"]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []

    def new_slide(title, subtitle=None, bg=None):
        slide = prs.slides.add_slide(blank)
        set_bg(slide, bg or COLORS["sand"])
        add_title(slide, title, subtitle)
        slides.append(slide)
        add_footer(slide, len(slides))
        return slide

    # 1
    s = prs.slides.add_slide(blank)
    set_bg(s, COLORS["sand"])
    title_box = s.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Design of a Plant to Manufacture Benzalkonium Chloride"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["navy"]
    p2 = tf.add_paragraph()
    p2.text = "(50 wt% Solution Basis)"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLORS["teal"]
    p3 = tf.add_paragraph()
    p3.text = "50,000 TPA | Continuous Process | Dahej, Gujarat Basis"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLORS["slate"]
    add_note_box(s, "Home Paper Presentation\nInstitute / Department\nStudent Name | Roll No.\nGuide Name", Inches(0.95), Inches(4.6), Inches(4.4), Inches(1.4), fill=COLORS["white"], font_size=15)
    add_footer(s, 1)
    slides.append(s)

    # 2
    s = new_slide("Presentation Roadmap")
    add_bullets(s, [
        "1. Product and commercial basis",
        "2. Route selection and process basis",
        "3. Technical feasibility and process flow",
        "4. Reactor and equipment design",
        "5. Control, safety, and plant layout",
        "6. Costing, working capital, and financial analysis",
    ], Inches(1.0), Inches(1.45), Inches(10.2), Inches(4.6), font_size=20)

    # 3
    s = new_slide("Benzalkonium Chloride: Introduction")
    add_bullets(s, [
        "Benzalkonium chloride (BAC) is a quaternary ammonium compound used as a biocide and cationic surfactant.",
        "It is used in disinfectants, sanitizers, cleaning formulations, and selected personal-care and industrial products.",
        "Commercial BAC is generally supplied as an active solution rather than as a neat pure compound.",
        "The present study is developed on a 50 wt% sold-solution basis for continuous manufacture.",
    ], Inches(0.9), Inches(1.5), Inches(11.3), Inches(4.8))

    # 4
    s = new_slide("Product Profile")
    add_table(s, 8, 2, [
        ["Basis Item", "Value"],
        ["Product", "Benzalkonium Chloride (BAC)"],
        ["Commercial form", "50 wt% sold-solution basis"],
        ["Throughput basis", "Finished product"],
        ["Active content", "50 wt%"],
        ["Sold solution rate", "6,313.13 kg/h"],
        ["Active rate", "3,156.57 kg/h"],
        ["Selling price basis", "INR 360.00/kg sold solution"],
    ], Inches(0.8), Inches(1.45), Inches(7.0), Inches(4.7), font_size=13)
    add_note_box(s, "Commercial BAC is treated here on a sold-solution basis rather than as a neat pure-component product.", Inches(8.25), Inches(1.7), Inches(4.0), Inches(1.0), font_size=13)

    # 5
    s = new_slide("Market and Capacity Justification")
    add_bullets(s, [
        "BAC has a strong market in disinfection, hygiene, and specialty chemical applications.",
        "A 50,000 TPA plant basis is selected to capture scale advantage while supporting domestic supply and export potential.",
        "The case is aligned with an India-based commercial manufacturing context.",
    ], Inches(0.8), Inches(1.55), Inches(6.0), Inches(4.2), font_size=17)
    add_table(s, 7, 2, [
        ["Parameter", "Value"],
        ["Installed capacity", "50,000 TPA"],
        ["Annual operating days", "330 d/y"],
        ["Product basis", "50 wt% BAC solution"],
        ["Throughput basis", "Finished product"],
        ["Site basis", "Dahej, Gujarat"],
        ["Economic basis year", "2025"],
    ], Inches(7.05), Inches(1.4), Inches(5.3), Inches(3.9), font_size=12)

    # 6
    s = new_slide("Literature and Route Options")
    add_table(s, 4, 5, [
        ["Route", "Family", "Score", "Status", "Key Note"],
        ["Solvent-Free (Neat) Quaternization", "Quaternization liquid train", "93.71", "Selected", "Best overall route score"],
        ["Quaternization in Butanone", "Quaternization liquid train", "72.94", "Not selected", "Lower overall score"],
        ["Quaternization in Acetonitrile", "Quaternization liquid train", "41.67", "Not selected", "Least favorable score"],
    ], Inches(0.55), Inches(1.45), Inches(12.2), Inches(3.0), font_size=11)
    add_note_box(s, "All identified routes belong to the same broad quaternization family, but differ in medium, operability burden, and economic screening outcome.", Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.8), font_size=13)

    # 7
    s = new_slide("Selected Process Basis")
    add_table(s, 7, 2, [
        ["Selection Item", "Value"],
        ["Selected route", "Solvent-Free (Neat) Quaternization"],
        ["Route family", "Quaternization liquid train"],
        ["Scientific status", "Screening-feasible"],
        ["Operating mode", "Continuous"],
        ["Main downstream train", "Primary separation -> concentration -> purification"],
        ["Selection caution", "Feasibility-grade academic basis; not a detailed industrial design package"],
    ], Inches(0.85), Inches(1.45), Inches(11.6), Inches(4.05), font_size=12)
    add_note_box(s, "The solvent-free quaternization route is retained as the current feasibility-grade design basis because it gave the strongest combined process and economic screening score under the present assumptions.", Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.82), font_size=12)

    # 8
    s = new_slide("Locked Process Basis")
    add_table(s, 9, 2, [
        ["Basis Item", "Value"],
        ["Product basis", "50 wt% BAC solution"],
        ["Capacity", "50,000 TPA"],
        ["Route family", "Quaternization liquid train"],
        ["Process train steps", "7"],
        ["Separation duties", "2"],
        ["Heat integration case", "No recovery"],
        ["Residual hot utility", "359.13 kW"],
        ["Residual cold utility", "567.12 kW"],
    ], Inches(1.05), Inches(1.55), Inches(10.7), Inches(4.4), font_size=13)

    # 9
    s = new_slide("Thermodynamic Feasibility")
    add_bullets(s, [
        "The selected BAC route is carried forward as a thermodynamically feasible liquid-phase quaternization system.",
        "The process basis supports reactor operation, cleanup, and finishing without requiring active-product distillation.",
        "Separation design is therefore driven primarily by volatile removal and purification rather than by BAC evaporation.",
        "The commercial sold-solution basis is consistent with the selected process architecture.",
    ], Inches(0.85), Inches(1.45), Inches(11.4), Inches(4.5))
    add_note_box(s, "Thermodynamic interpretation is developed on the current feasibility-grade pseudo-component and cleanup basis.", Inches(0.9), Inches(5.7), Inches(11.3), Inches(0.7), font_size=12)

    # 10
    s = new_slide("Reaction Kinetics and Design Basis")
    add_bullets(s, [
        "The quaternization step is treated on a screening kinetic basis suitable for preliminary reactor sizing.",
        "The selected design case assumes high conversion with a controlled exothermic liquid-phase reaction environment.",
        "Reactor sizing is carried forward on a 25 h residence-time basis with 95% design conversion.",
        "The kinetic basis is sufficient for feasibility-stage design, but still requires refinement for detailed industrial validation.",
    ], Inches(0.85), Inches(1.45), Inches(8.1), Inches(4.6), font_size=16)
    add_metric_cards(s, [("Residence time", "25.0 h"), ("Design conversion", "0.95")], Inches(9.35), Inches(1.85), Inches(2.4), Inches(1.0))

    # 11
    s = new_slide("Block Flow Diagram")
    add_image_contain(s, ASSET_DIR / "bfd_sheet_1.png", Inches(0.6), Inches(1.2), Inches(8.0), Inches(5.5))
    add_bullets(s, [
        "Feed preparation and charging",
        "Continuous quaternization reactor",
        "Primary flash, concentration, and purification",
        "Storage, dispatch, and waste handling",
    ], Inches(8.85), Inches(1.65), Inches(3.7), Inches(3.8), font_size=15)

    # 12
    s = new_slide("Process Flow Diagram I")
    add_image_contain(s, ASSET_DIR / "pfd_sheet_1.png", Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.7))
    add_note_box(s, "Feed preparation and reaction section of the selected BAC process train.", Inches(3.6), Inches(6.45), Inches(6.2), Inches(0.45), font_size=11)

    # 13
    s = new_slide("Process Flow Diagram II")
    add_image_contain(s, ASSET_DIR / "pfd_sheet_2.png", Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.7))
    add_note_box(s, "Cleanup, finishing, storage, and waste-handling section of the selected BAC process train.", Inches(3.2), Inches(6.45), Inches(7.0), Inches(0.45), font_size=11)

    # 14
    s = new_slide("Process Description")
    add_bullets(s, [
        "1. Feed components are received, stored, and prepared in the feed-handling section.",
        "2. Reactants enter the continuous quaternization reactor, where BAC is formed in the liquid phase.",
        "3. Reactor effluent passes to primary flash for volatile removal and downstream stabilization.",
        "4. The partially cleaned stream is concentrated and stripped to reduce volatile and light-end burden.",
        "5. A purification section refines the product to the required sold-solution quality basis.",
        "6. Finished BAC solution is transferred to storage; purge and waste streams are routed to controlled handling.",
    ], Inches(0.85), Inches(1.35), Inches(11.5), Inches(4.9), font_size=16)

    # 15
    s = new_slide("Material Balance Summary")
    add_table(s, 7, 2, [
        ["Stream / Basis Item", "Value"],
        ["Feed to reactor", "3,461.60 kg/h"],
        ["Product basis", "6,313.13 kg/h sold solution"],
        ["Active basis", "3,156.57 kg/h"],
        ["Design conversion", "95.0%"],
        ["Major feed components", "Dodecyldimethylamine, benzyl chloride"],
        ["Major product basis", "BAC solution with cleanup and finishing"],
    ], Inches(0.85), Inches(1.45), Inches(11.5), Inches(4.2), font_size=13)
    add_note_box(s, "This slide presents the plant balance basis only; the full stream ledger is intentionally kept out of the main deck.", Inches(0.9), Inches(5.85), Inches(11.3), Inches(0.65), font_size=12)

    # 16
    s = new_slide("Energy Balance Summary")
    add_table(s, 6, 4, [
        ["Utility", "Load", "Units", "Annualized Usage"],
        ["Steam", "2,277.14", "kg/h", "18,034,917.1 kg/y"],
        ["Cooling water", "107.59", "m3/h", "852,136.6 m3/y"],
        ["Electricity", "151.56", "kW", "1,200,323.5 kWh/y"],
        ["DM water", "2.00", "m3/h", "15,840.0 m3/y"],
        ["Nitrogen", "3.00", "Nm3/h", "23,760.0 Nm3/y"],
    ], Inches(0.75), Inches(1.45), Inches(12.0), Inches(3.7), font_size=12)
    add_note_box(s, "The thermal burden is dominated by steam and cooling-water service linked to concentration, purification, and reactor heat removal.", Inches(0.85), Inches(5.45), Inches(11.6), Inches(0.72), font_size=12)

    # 17
    s = new_slide("Heat Integration and Utilities Basis")
    add_bullets(s, [
        "The current BAC design is carried forward on a no-recovery heat-integration case.",
        "Residual utility demand remains at 359.13 kW hot utility and 567.12 kW cold utility.",
        "Steam, cooling water, electricity, nitrogen, and DM water are the principal utility services.",
        "Utility routing and control are integrated with the selected continuous liquid-train process architecture.",
    ], Inches(0.85), Inches(1.5), Inches(7.9), Inches(4.4), font_size=16)
    add_metric_cards(s, [("Hot utility", "359.13 kW"), ("Cold utility", "567.12 kW")], Inches(8.95), Inches(1.95), Inches(2.55), Inches(1.0))

    # 18
    s = new_slide("Reactor Design Basis")
    add_table(s, 9, 2, [
        ["Parameter", "Value"],
        ["Reactor tag", "R-101"],
        ["Adopted reactor basis", "Glass-lined agitated reactor with external recirculation loop"],
        ["Residence time", "25.0 h"],
        ["Design conversion", "0.95"],
        ["Design temperature", "85.0 degC"],
        ["Design pressure", "3.01 bar"],
        ["Heat duty", "257.12 kW"],
        ["Runaway risk label", "Moderate"],
    ], Inches(0.8), Inches(1.35), Inches(7.3), Inches(4.8), font_size=12)
    add_note_box(s, "R-101 is presented on the adopted glass-lined, externally cooled reactor basis selected for safe preliminary handling of the exothermic quaternization service.", Inches(8.35), Inches(1.8), Inches(4.1), Inches(1.5), font_size=13)

    # 19
    s = new_slide("Reactor Sizing Summary")
    add_table(s, 8, 2, [
        ["Parameter", "Value"],
        ["Design volume", "123.76 m3"],
        ["Liquid holdup", "101.44 m3"],
        ["Shell diameter", "4.00 m"],
        ["Shell length", "9.85 m"],
        ["Vessel height without skirt", "11.85 m"],
        ["External-loop exchanger area", "15.00 m2"],
        ["Agitator arrangement", "3 x retreat-curve agitator"],
    ], Inches(0.85), Inches(1.4), Inches(6.7), Inches(4.5), font_size=13)
    add_note_box(s, "The adopted geometry reflects the final report-basis reactor sizing and excludes superseded jacket-only wording from historical artifacts.", Inches(7.95), Inches(1.95), Inches(4.2), Inches(1.25), font_size=13)

    # 20
    s = new_slide("Reactor Mechanical Design and MoC")
    add_table(s, 7, 2, [
        ["Parameter", "Value"],
        ["Reactor type", "Glass-lined vertical reactor"],
        ["Pressure boundary basis", "Carbon-steel shell with glass lining"],
        ["Head type", "2:1 ellipsoidal heads"],
        ["Cooling philosophy", "External pump-around loop with exchanger"],
        ["Heat-transfer area", "15.00 m2"],
        ["Runaway risk category", "Moderate"],
    ], Inches(0.8), Inches(1.45), Inches(6.8), Inches(3.9), font_size=13)
    add_bullets(s, [
        "Glass lining is preferred to support corrosive quaternary-ammonium service.",
        "The external recirculation loop provides controlled heat removal for the exothermic reaction.",
        "Mechanical design remains preliminary and should be refined further during detailed design.",
    ], Inches(7.95), Inches(1.7), Inches(4.0), Inches(3.3), font_size=15)

    # 21
    s = new_slide("Major Process Unit Design")
    add_table(s, 5, 7, [
        ["ID", "Unit", "Service", "Volume", "Duty", "Design T", "Design P"],
        ["R-101", "Reactor", "Quaternization reaction", "123.76 m3", "257.12 kW", "85.0 degC", "3.01 bar"],
        ["PU-201", "Distillation column", "Purification train", "21.21 m3", "1,391.58 kW", "140.0 degC", "2.00 bar"],
        ["V-101", "Flash drum", "Intermediate disengagement", "22.28 m3", "0.00 kW", "85.0 degC", "3.00 bar"],
        ["E-101", "Heat exchanger", "Thermal service", "8.17 m3", "637.18 kW", "180.0 degC", "8.00 bar"],
    ], Inches(0.35), Inches(1.35), Inches(12.6), Inches(3.4), font_size=10)
    add_note_box(s, "The selected BAC process is carried by a compact unit train dominated by reaction, volatile cleanup, purification, and storage service.", Inches(0.8), Inches(5.15), Inches(11.5), Inches(0.7), font_size=12)

    # 22
    s = new_slide("Heat Exchanger and Thermal Equipment")
    add_table(s, 3, 6, [
        ["Equipment", "Service", "Duty", "Design T", "Design P", "MoC"],
        ["E-101", "Shell-and-tube exchanger", "637.18 kW", "180.0 degC", "8.00 bar", "Carbon steel"],
        ["R-101 external loop", "Reactor heat removal", "257.12 kW", "85.0 degC", "3.01 bar", "Integrated with GLR loop"],
    ], Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.3), font_size=12)
    add_note_box(s, "Thermal equipment design is driven by reactor temperature control and volatile-cleanup / purification duties.", Inches(0.9), Inches(4.45), Inches(11.2), Inches(0.75), font_size=12)

    # 23
    s = new_slide("Pumps, Storage, and Auxiliary Equipment")
    add_table(s, 3, 4, [
        ["Equipment", "Service", "Key Size / Duty", "Design Basis"],
        ["TK-301", "BAC product storage tank", "1,969.85 m3 total volume", "11.6 days inventory"],
        ["TK-301 Pump", "Product transfer", "219.85 m3/h, 52.0 m head, 36.58 kW", "Dispatch and tank-farm transfer"],
    ], Inches(0.75), Inches(1.65), Inches(12.0), Inches(2.3), font_size=12)
    add_note_box(s, "Storage and transfer equipment are sized to support continuous production, dispatch buffering, and inventory control under the selected operating policy.", Inches(0.9), Inches(4.45), Inches(11.2), Inches(0.8), font_size=12)

    # 24
    s = new_slide("Equipment Design Summary")
    add_bullets(s, [
        "The BAC plant is built around a continuous quaternization and cleanup train with a limited number of major equipment items.",
        "Reactor R-101 remains the key safety-critical and design-critical unit.",
        "PU-201 and E-101 dominate the thermal and purification burden.",
        "TK-301 and its transfer service define the product-handling and dispatch basis.",
        "The overall equipment set is consistent with a compact continuous liquid-process facility.",
    ], Inches(0.85), Inches(1.5), Inches(11.4), Inches(4.8), font_size=17)

    # 25
    s = new_slide("Instrumentation and Control Philosophy")
    add_table(s, 6, 2, [
        ["Parameter", "Value"],
        ["Selected architecture", "SIS-augmented regulatory control"],
        ["Critical units", "R-101, PU-201, TK-301"],
        ["Loop count", "8"],
        ["Utility services linked", "6"],
        ["High-criticality loops", "4"],
    ], Inches(0.8), Inches(1.5), Inches(6.5), Inches(3.4), font_size=13)
    add_bullets(s, [
        "Reactor temperature, pressure, and feed ratio form the critical control core.",
        "Purification pressure and inventory control protect product quality and operability.",
        "Storage blanketing and level control support safe dispatch and containment.",
    ], Inches(7.85), Inches(1.8), Inches(4.0), Inches(3.4), font_size=15)

    # 26
    s = new_slide("Control System Diagram")
    add_image_contain(s, ASSET_DIR / "control_system_sheet_1.png", Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.7))
    add_note_box(s, "Principal BAC control loops centered on reactor safety, purification stability, and storage handling.", Inches(3.2), Inches(6.45), Inches(7.0), Inches(0.45), font_size=11)

    # 27
    s = new_slide("HAZOP Overview")
    add_table(s, 5, 2, [
        ["Parameter", "Value"],
        ["Node count", "4"],
        ["High-severity route hazards", "3"],
        ["Control loops linked", "8"],
        ["Node families", "Reactor, separation, thermal exchange, storage"],
    ], Inches(1.2), Inches(1.7), Inches(5.6), Inches(2.8), font_size=13)
    add_note_box(s, "The HAZOP is focused on the major BAC risk nodes rather than on an exhaustive detailed-plant register.", Inches(7.25), Inches(2.1), Inches(4.7), Inches(1.0), font_size=13)

    # 28
    s = new_slide("Critical HAZOP Deviations")
    add_table(s, 5, 5, [
        ["Node", "Deviation", "Severity", "Main Consequence", "Safeguards"],
        ["R-101", "High reactor temperature", "High", "Runaway tendency, selectivity loss, pressure rise", "TIC-R-101; PIC-R-101; FRC-R-101"],
        ["PU-201", "Low separation pressure", "High", "Air ingress, instability, off-spec separation", "LIC-PU-201; PIC-PU-201; FIC-PU-201"],
        ["E-101", "Low thermal duty", "Medium", "Poor temperature control and downstream upset", "TIC-R-101; PIC-R-101"],
        ["TK-301", "High storage level", "Medium", "Overflow and containment loss", "LIC-TK-301; PIC-TK-301"],
    ], Inches(0.35), Inches(1.35), Inches(12.6), Inches(3.5), font_size=10)
    add_note_box(s, "Reactor temperature control and storage containment remain the most important safety anchors in the present BAC process basis.", Inches(0.85), Inches(5.1), Inches(11.3), Inches(0.78), font_size=12)

    # 29
    s = new_slide("SHE, Waste, and ETP Basis")
    add_shew_diagram(s)
    add_note_box(s, "The SHE basis is aligned with hazardous-chemical handling, segregated waste management, and a no-direct-discharge operating philosophy.", Inches(0.9), Inches(5.35), Inches(11.2), Inches(0.72), font_size=12)

    # 30
    s = new_slide("Site Selection")
    add_table(s, 7, 2, [
        ["Parameter", "Value"],
        ["Selected site", "Dahej, Gujarat"],
        ["Region", "India"],
        ["Industrial basis", "Coastal chemical cluster / PCPIR-linked basis"],
        ["Utility basis", "Grounded in cited site utility basis"],
        ["Logistics basis", "Coastal cluster dispatch basis"],
        ["Site-fit note", "Best available fit for BAC continuous liquid chemical manufacture"],
    ], Inches(0.8), Inches(1.45), Inches(6.9), Inches(4.0), font_size=12)
    add_bullets(s, [
        "Strong feedstock ecosystem and chemical-cluster support",
        "Good coastal and land logistics",
        "Suitable utility and regulatory basis for hazardous liquid chemical manufacture",
    ], Inches(8.0), Inches(1.9), Inches(4.0), Inches(2.8), font_size=15)

    # 31
    s = new_slide("Plant Layout Basis")
    add_zone_diagram(s)
    add_bullets(s, [
        "Tank farm and receipt zone positioned near feed entry and dispatch edge",
        "Reaction zone segregated from occupied and dispatch areas",
        "Separation and finishing zone located downstream of reaction for short process transfer",
        "Waste / recovery zone isolated but connected to the main process train",
    ], Inches(0.9), Inches(5.35), Inches(11.3), Inches(1.2), font_size=13)

    # 32
    s = new_slide("Layout and Utility Corridors")
    add_utility_diagram(s)
    add_bullets(s, [
        "Steam and cooling-water headers run parallel to the process train with branch take-offs at major units.",
        "Electrical and service corridors are routed away from wet and hot process hazard zones.",
        "Nitrogen service is aligned with storage blanketing and inerting requirements.",
        "Truck movement, utility routing, emergency response, and maintenance access are kept separate wherever practical.",
    ], Inches(0.8), Inches(4.25), Inches(11.6), Inches(2.0), font_size=13)

    # 33
    s = new_slide("Equipment Costing and Project Cost")
    add_table(s, 7, 2, [
        ["CAPEX Head", "Value"],
        ["Purchased equipment", "34.45 Cr INR"],
        ["Installed equipment", "81.00 Cr INR"],
        ["Direct plant cost", "93.98 Cr INR"],
        ["Indirect cost", "22.56 Cr INR"],
        ["Contingency", "9.40 Cr INR"],
        ["Total CAPEX", "125.93 Cr INR"],
    ], Inches(0.65), Inches(1.35), Inches(6.2), Inches(4.1), font_size=12)
    add_table(s, 6, 3, [
        ["Equipment Family", "Count", "Installed Cost"],
        ["Storage tank", "1", "47.36 Cr INR"],
        ["Reactor", "1", "19.99 Cr INR"],
        ["Distillation column", "1", "2.23 Cr INR"],
        ["Heat exchanger", "1", "1.05 Cr INR"],
        ["Flash drum", "1", "0.97 Cr INR"],
    ], Inches(7.15), Inches(1.55), Inches(5.3), Inches(3.0), font_size=11)
    add_note_box(s, "Project cost has been estimated on a screening basis to support feasibility assessment; further refinement is required during detailed design.", Inches(0.85), Inches(5.55), Inches(11.4), Inches(0.72), font_size=12)

    # 34
    s = new_slide("Cost of Production")
    add_table(s, 7, 3, [
        ["OPEX Head", "Value", "Share of OPEX"],
        ["Raw materials", "586.59 Cr INR/y", "95.49%"],
        ["Utilities", "4.95 Cr INR/y", "0.81%"],
        ["Labor", "15.60 Cr INR/y", "2.54%"],
        ["Maintenance", "3.54 Cr INR/y", "0.58%"],
        ["Overheads", "3.64 Cr INR/y", "0.59%"],
        ["Total OPEX", "614.31 Cr INR/y", "100.00%"],
    ], Inches(0.7), Inches(1.45), Inches(12.0), Inches(3.5), font_size=12)
    add_note_box(s, "The BAC case is strongly raw-material-cost dominated, while utility and maintenance burdens remain comparatively smaller at the present feasibility basis.", Inches(0.85), Inches(5.3), Inches(11.4), Inches(0.75), font_size=12)

    # 35
    s = new_slide("Working Capital and Financial Analysis")
    add_table(s, 7, 2, [
        ["Component", "Value"],
        ["Raw-material inventory", "27.12 Cr INR"],
        ["Product inventory", "41.92 Cr INR"],
        ["Receivables", "157.81 Cr INR"],
        ["Cash buffer", "11.78 Cr INR"],
        ["Payables", "38.57 Cr INR"],
        ["Working capital", "204.62 Cr INR"],
    ], Inches(0.6), Inches(1.4), Inches(5.6), Inches(3.7), font_size=12)
    add_metric_cards(s, [
        ("Revenue", "1,800.00 Cr INR/y"),
        ("OPEX", "614.31 Cr INR/y"),
        ("Payback", "0.48 y"),
        ("IRR", "60.00%"),
    ], Inches(6.7), Inches(1.75), Inches(1.4), Inches(1.0), gap=0.15)
    add_note_box(s, "NPV: 4,366.57 Cr INR | Break-even capacity: 34.10%", Inches(6.75), Inches(3.3), Inches(5.2), Inches(0.65), fill=RGBColor(247, 248, 252), font_size=13)
    add_note_box(s, "These financial indicators are model-based feasibility outputs and remain sensitive to route assumptions, purification performance, and selling-price basis.", Inches(0.85), Inches(5.55), Inches(11.4), Inches(0.78), fill=RGBColor(250, 244, 232), font_size=12)

    # 36
    s = new_slide("Conclusion")
    add_bullets(s, [
        "A continuous BAC plant on a 50 wt% sold-solution basis has been developed on a coherent feasibility-grade design basis.",
        "The selected quaternization route, process flow, and unit train are technically consistent with the present commercial product basis.",
        "Reactor, purification, storage, control, and safety sections have been carried through preliminary design and risk review.",
        "Dahej, Gujarat provides a suitable site basis for utilities, logistics, and hazardous-chemical industrial support.",
        "The current economic results indicate a promising case under the adopted assumptions, but require further refinement in route detail, purification basis, and project costing.",
        "Overall, the BAC plant appears technically coherent and economically promising on the current feasibility basis.",
    ], Inches(0.85), Inches(1.45), Inches(11.4), Inches(4.9), font_size=17)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"saved {OUT_PATH}")


if __name__ == "__main__":
    build()
